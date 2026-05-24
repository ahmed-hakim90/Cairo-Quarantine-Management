#!/usr/bin/env python3
"""Generate Arabic RTL ministry presentation PPTX for Cairo Quarantine portal."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

import arabic_reshaper
from bidi.algorithm import get_display
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
OUT_PATH = DOCS / "Cairo-Quarantine-Ministry.pptx"
SHOTS = DOCS / "handover" / "screenshots"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0x1A, 0x2B, 0x3C)
MID = RGBColor(0x5A, 0x6B, 0x7D)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)

FONT = "Arial Unicode MS"
DEVELOPERS = [
    ("مطور ١", "[الجهة / الشركة]"),
    ("مطور ٢", "[الجهة / الشركة]"),
]


def ar(text: str) -> str:
    if not text:
        return ""
    return get_display(arabic_reshaper.reshape(text))


@dataclass
class SlideSpec:
    title: str
    bullets: list[str] = field(default_factory=list)
    body: str = ""
    images: list[str] = field(default_factory=list)
    table: list[list[str]] | None = None
    layout: str = "content"  # content | cover | two_images | table


def shot(name: str) -> Path | None:
    path = SHOTS / name
    return path if path.is_file() else None


def fit_image(path: Path, max_w: float, max_h: float) -> tuple[float, float]:
    with PILImage.open(path) as im:
        w, h = im.size
    ratio = min(max_w / w, max_h / h)
    return w * ratio, h * ratio


def set_rtl_paragraph(paragraph, *, size: int = 18, bold: bool = False, color: RGBColor = BODY) -> None:
    paragraph.alignment = PP_ALIGN.RIGHT
    for run in paragraph.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = BODY,
    align: PP_ALIGN = PP_ALIGN.RIGHT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = ar(text)
    p.alignment = align
    if p.runs:
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.bold = bold
        p.runs[0].font.color.rgb = color


def add_bullets(slide, left: float, top: float, width: float, height: float, items: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ar(f"• {item}")
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(8)
        if p.runs:
            p.runs[0].font.name = FONT
            p.runs[0].font.size = Pt(16)
            p.runs[0].font.color.rgb = BODY


def add_header(slide, title: str) -> None:
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        SLIDE_W,
        Inches(0.85),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0.82), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    add_textbox(slide, 0.4, 0.12, 12.5, 0.55, title, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        0.4,
        0.52,
        12.5,
        0.25,
        "بوابة إدارة الحجر الصحي بالقاهرة",
        size=10,
        color=RGBColor(0xD8, 0xE2, 0xF0),
        align=PP_ALIGN.CENTER,
    )


def add_image(slide, filename: str, left: float, top: float, max_w: float, max_h: float) -> bool:
    path = shot(filename)
    if not path:
        return False
    w, h = fit_image(path, max_w * 96, max_h * 96)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(w / 96), height=Inches(h / 96))
    return True


def add_table(slide, rows: list[list[str]], left: float, top: float, width: float) -> None:
    table_rows = len(rows)
    table_cols = len(rows[0]) if rows else 0
    shape = slide.shapes.add_table(table_rows, table_cols, Inches(left), Inches(top), Inches(width), Inches(0.35 * table_rows))
    table = shape.table
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ar(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.RIGHT
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(12 if r else 13)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else BODY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT


def build_slides() -> list[SlideSpec]:
    return [
        SlideSpec(
            title="بوابة إدارة الحجر الصحي بالقاهرة",
            layout="cover",
            body="عرض تقديمي لوزارة الصحة\nمايو 2026",
        ),
        SlideSpec(
            title="مقدمة",
            bullets=[
                "نقدّم منصة رقمية موحّدة لخدمات الحجر الصحي والتطعيمات بالقاهرة.",
                "المنصة تجمع المعلومة، الحجز، الشكاوى، ومتابعة الدور في مكان واحد.",
                "الهدف: تسهيل وصول المواطن للخدمة وتنظيم عمل المكاتب.",
            ],
        ),
        SlideSpec(
            title="المشكلة",
            bullets=[
                "كان المواطن يبحث عن المعلومة في أكثر من مصدر أو يسأل في المكتب.",
                "الحجز كان يتطلب الحضور الشخصي للمكتب.",
                "معرفة الدور والانتظار لم تكن واضحة أو قابلة للمتابعة من البيت.",
            ],
        ),
        SlideSpec(
            title="الحل — ثلاث نقاط أساسية",
            bullets=[
                "① معلومة واضحة ومتاحة على مدار الساعة لكل أنواع السفر.",
                "② حجز موعد من المنزل مع بطاقة رقمية (QR) للعرض عند المكتب.",
                "③ متابعة الطلب والدور رقمياً دون الحاجة للاتصال المتكرر.",
            ],
            images=["03-public-booking-ar.png"],
        ),
        SlideSpec(
            title="الصفحة الرئيسية",
            bullets=[
                "بوابة واحدة تجمع الإرشاد، جدول اللقاحات، ومواقع المكاتب.",
                "روابط سريعة لكل فئات المسافرين وخدمات المواطن.",
            ],
            images=["01-public-home-ar.png"],
        ),
        SlideSpec(
            title="أنواع السفر",
            bullets=[
                "مسافر دولي — متطلبات التطعيم حسب وجهة السفر.",
                "حج وعمرة — إرشادات ومكاتب وأسعار مخصّصة.",
                "مواطن — خدمات التطعيم المحلية والإرشاد الصحي.",
                "كل فئة لها صفحة إرشادية تفاعلية داخل المنصة.",
            ],
            images=[
                "02-public-international-traveler-ar.png",
                "07-public-hajj-umrah-ar.png",
                "08-public-citizen-services-ar.png",
            ],
            layout="three_images",
        ),
        SlideSpec(
            title="مسافر دولي — متطلبات الدولة",
            bullets=[
                "بيئة تفاعلية: يختار المواطن الدولة فتظهر التطعيمات المطلوبة فوراً.",
                "البيانات تُحدَّث من لوحة الإدارة دون تعديل البرمجيات.",
            ],
            images=["02-public-international-traveler-ar.png"],
        ),
        SlideSpec(
            title="مثال: المملكة العربية السعودية",
            bullets=[
                "عند اختيار المملكة العربية السعودية تظهر قائمة التطعيمات والاشتراطات المطلوبة.",
                "يساعد المسافر على التحضير قبل زيارة المكتب.",
                "يمكن إضافة أو تعديل متطلبات أي دولة من لوحة السوبر أدمن.",
            ],
            images=["saudi-requirements-ar.png"],
        ),
        SlideSpec(
            title="حج وعمرة",
            bullets=[
                "صفحة مخصّصة للحجاج والمعتمرين: إرشادات صحية وجدول المكاتب.",
                "عرض الأسعار والمواقع لاتخاذ القرار قبل الحجز.",
            ],
            images=["07-public-hajj-umrah-ar.png"],
        ),
        SlideSpec(
            title="خدمات المواطن",
            bullets=[
                "إرشادات التطعيم للمواطنين داخل مصر.",
                "معلومات واضحة تساعد على معرفة ما يحتاجه قبل زيارة المكتب.",
            ],
            images=["08-public-citizen-services-ar.png"],
        ),
        SlideSpec(
            title="حجز موعد",
            bullets=[
                "نموذج بسيط: حالة المسافر، المكتب، التاريخ، الاسم والهاتف.",
                "النظام يتحقق من السعة اليومية للمكتب قبل تأكيد الحجز.",
                "بعد الإرسال يحصل المواطن على رقم طلب فوري.",
            ],
            images=["03-public-booking-ar.png"],
        ),
        SlideSpec(
            title="بطاقة الحجز (QR)",
            bullets=[
                "بطاقة إلكترونية برمز QR للعرض عند المكتب.",
                "يمكن حفظها أو مشاركتها — صالحة لمدة محددة.",
                "تسهّل على الموظف التحقق من الحجز بسرعة.",
            ],
            images=["14-public-booking-pass-ar.png"],
        ),
        SlideSpec(
            title="الطابور اليومي",
            bullets=[
                "عند الوصول للمكتب: مسح QR الحضور أو إدخال رقم الطلب.",
                "المواطن يعرف دوره ووقت الانتظار المتوقع.",
                "الموظف يدير الطابور من شاشة مخصّصة في المكتب.",
            ],
            images=["06-public-checkin-ar.png", "19-admin-office-queue-ar.png"],
            layout="two_images",
        ),
        SlideSpec(
            title="الشكاوى — جانب المواطن",
            bullets=[
                "نموذج إلكتروني للشكوى أو المقترح مرتبط بمكتب محدد.",
                "تفاصيل إلزامية لضمان وصول الشكوى للجهة المعنية.",
                "متابعة الحالة من صفحة «طلباتي» على الجهاز.",
            ],
            images=["04-public-complaint-ar.png", "05-public-my-requests-ar.png"],
            layout="two_images",
        ),
        SlideSpec(
            title="الشكاوى — جانب الموظف",
            bullets=[
                "① وصول الشكوى إلى قائمة الطلبات في لوحة الإدارة.",
                "② التواصل مع المواطن عبر واتساب من قالب جاهز.",
                "③ تحديث الحالة: جديد → تم التواصل → قيد المعالجة → مكتمل.",
                "④ سجل نشاط كامل للتدقيق والمتابعة.",
            ],
            images=["12-admin-requests-ar.png", "13-admin-request-detail-ar.png"],
            layout="two_images",
        ),
        SlideSpec(
            title="إدارة المكاتب",
            bullets=[
                "إضافة وتعديل المكاتب: العنوان، الهاتف، ساعات العمل.",
                "تحديد السعة اليومية للحجز لكل مكتب.",
                "إنشاء QR حضور خاص بكل مكتب للطابور اليومي.",
            ],
            images=["15-admin-offices-ar.png"],
        ),
        SlideSpec(
            title="إدارة التطعيمات والبيانات",
            bullets=[
                "السوبر أدمن يتحكم في: اللقاحات، متطلبات الدول، حالات المسافرين.",
                "تحديث البيانات من لوحة واحدة أو استيراد Excel.",
                "لا حاجة لمطوّر لتغيير الأسعار أو إضافة دولة جديدة.",
            ],
            images=["20-admin-vaccines-ar.png", "16-admin-destination-countries-ar.png"],
            layout="two_images",
        ),
        SlideSpec(
            title="فريق التطوير",
            body="أفراد مستقلون يعملون على تقديم حل تقني لخدمة المواطن وتسهيل عمل المكاتب.",
            layout="developers",
        ),
        SlideSpec(
            title="التقنيات الحالية",
            table=[
                ["البند", "التفصيل"],
                ["نوع المنصة", "موقع ويب حديث يعمل على الجوال والكمبيوتر"],
                ["اللغات", "عربي (افتراضي) · إنجليزي · صيني"],
                ["قاعدة البيانات", "خدمة سحابية آمنة (Firebase)"],
                ["الأمان", "دخول الموظفين محمي · حدود للطلبات المتكررة"],
            ],
            layout="table",
        ),
        SlideSpec(
            title="خطة التوسع",
            bullets=[
                "المرحلة 1 (الآن): تشغيل مكاتب القاهرة — حتى ~50 ألف عملية يومياً.",
                "المرحلة 2 (قريباً): شبكة توزيع (CDN) + مراقبة + إشعارات للطابور.",
                "المرحلة 3 (جمهورية): خوادم متعددة + قاعدة بيانات أقوى + تكامل هوية وطنية.",
            ],
        ),
        SlideSpec(
            title="ما نحتاجه من الوزارة",
            bullets=[
                "بيئة تشغيل: استضافة سحابية + نطاق رسمي (.eg).",
                "حساب Firebase / Google Cloud للإنتاج (قاعدة بيانات + دخول الموظفين).",
                "جدولة مهام يومية: إغلاق الطابور، أرشفة، نسخ احتياطي.",
                "بيانات أولية: المكاتب، اللقاحات، متطلبات الدول.",
                "تدريب 1–2 يوم للسوبر أدمن ومديري المكاتب.",
                "إذا وفرت الوزارة البيئة — نبدأ التشغيل فوراً. وإلا نسعى لتوفير الأساسيات.",
            ],
        ),
        SlideSpec(
            title="شكراً لحضرتكم",
            layout="cover",
            body="نرحب بأسئلتكم وملاحظاتكم\nبوابة إدارة الحجر الصحي بالقاهرة",
        ),
    ]


def render_cover(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = slide.shapes.add_shape(1, Inches(0), Inches(3.0), SLIDE_W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    add_textbox(slide, 0.8, 2.2, 11.7, 1.2, spec.title, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if spec.body:
        add_textbox(slide, 0.8, 3.5, 11.7, 1.5, spec.body, size=18, color=RGBColor(0xB8, 0xC8, 0xDC), align=PP_ALIGN.CENTER)


def render_content(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, spec.title)

    has_images = bool(spec.images)
    text_top = 1.05
    text_w = 12.5
    text_h = 2.2 if has_images else 5.5

    if spec.body:
        add_textbox(slide, 0.4, text_top, text_w, 0.8, spec.body, size=17, color=BODY)
        text_top += 0.85

    if spec.bullets:
        add_bullets(slide, 0.4, text_top, text_w, text_h, spec.bullets)

    if not spec.images:
        return

    if spec.layout == "three_images" and len(spec.images) >= 3:
        y = 4.0
        w = 3.9
        for i, img in enumerate(spec.images[:3]):
            add_image(slide, img, 0.35 + i * 4.15, y, w, 3.0)
        return

    if spec.layout == "two_images" and len(spec.images) >= 2:
        add_image(slide, spec.images[0], 0.35, 3.85, 6.1, 3.2)
        add_image(slide, spec.images[1], 6.85, 3.85, 6.1, 3.2)
        return

    add_image(slide, spec.images[0], 3.2, 3.5, 6.8, 3.5)


def render_developers(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, spec.title)
    add_textbox(slide, 0.4, 1.1, 12.5, 0.8, spec.body, size=17, color=BODY)
    y = 2.4
    for name, org in DEVELOPERS:
        card = slide.shapes.add_shape(1, Inches(2.5), Inches(y), Inches(8.3), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = ACCENT
        add_textbox(slide, 2.7, y + 0.15, 7.9, 0.45, name, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, 2.7, y + 0.58, 7.9, 0.4, org, size=14, color=MID, align=PP_ALIGN.CENTER)
        y += 1.45


def render_table(prs: Presentation, spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, spec.title)
    if spec.table:
        add_table(slide, spec.table, 1.2, 1.5, 10.9)


def build() -> None:
    SHOTS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for spec in build_slides():
        if spec.layout == "cover":
            render_cover(prs, spec)
        elif spec.layout == "developers":
            render_developers(prs, spec)
        elif spec.layout == "table":
            render_table(prs, spec)
        else:
            render_content(prs, spec)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    build()
